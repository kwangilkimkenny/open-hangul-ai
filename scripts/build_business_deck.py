"""Generate HanView business plan PPTX (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- Theme ----
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
INK = RGBColor(0x1A, 0x1F, 0x2C)
SLATE = RGBColor(0x4A, 0x55, 0x68)
MIST = RGBColor(0xE8, 0xEC, 0xF2)
LINE = RGBColor(0xD0, 0xD6, 0xDF)
ACCENT = RGBColor(0xE1, 0x34, 0x3B)  # HanView red
GOLD = RGBColor(0xC9, 0x9A, 0x3B)
GREEN = RGBColor(0x2F, 0x7D, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF8, 0xFB)

FONT_TITLE = "맑은 고딕"
FONT_BODY = "맑은 고딕"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def header(slide, idx, total, title, subtitle=None, section=None):
    # top band
    add_rect(slide, 0, 0, SW, Inches(0.55), NAVY)
    add_text(slide, Inches(0.5), Inches(0.08), Inches(6), Inches(0.4),
             "HanView · OpenHangul-AI · 사업화 IR Deck",
             size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(10.5), Inches(0.08), Inches(2.7), Inches(0.4),
             f"{idx:02d} / {total:02d}", size=11, color=MIST,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # title block
    if section:
        add_text(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.3),
                 section, size=11, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=NAVY, font=FONT_TITLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.65), Inches(12), Inches(0.4),
                 subtitle, size=13, color=SLATE)
    # divider
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(0.5), Emu(25000), ACCENT)
    # footer
    add_rect(slide, 0, SH - Inches(0.3), SW, Inches(0.3), MIST)
    add_text(slide, Inches(0.5), SH - Inches(0.28), Inches(12.3), Inches(0.26),
             "Confidential — HanView 사업화 전략 · 2026Q2",
             size=9, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)


def kpi_card(slide, x, y, w, h, value, label, accent=ACCENT):
    add_rect(slide, x, y, w, h, WHITE, line=LINE)
    add_rect(slide, x, y, Inches(0.08), h, accent)
    add_text(slide, x + Inches(0.25), y + Inches(0.2), w - Inches(0.3), Inches(0.6),
             value, size=22, bold=True, color=NAVY, font=FONT_TITLE)
    add_text(slide, x + Inches(0.25), y + Inches(0.85), w - Inches(0.3), h - Inches(0.9),
             label, size=10, color=SLATE)


def bullet_box(slide, x, y, w, h, title, items, accent=ACCENT):
    add_rect(slide, x, y, w, h, WHITE, line=LINE)
    add_rect(slide, x, y, w, Inches(0.05), accent)
    add_text(slide, x + Inches(0.25), y + Inches(0.18), w - Inches(0.3), Inches(0.4),
             title, size=13, bold=True, color=NAVY)
    tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.65),
                                  w - Inches(0.4), h - Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = "• " + it
        r.font.name = FONT_BODY
        r.font.size = Pt(10.5)
        r.font.color.rgb = INK


def table(slide, x, y, w, h, headers, rows, col_widths=None, header_color=NAVY,
          zebra=True, font_size=10):
    cols = len(headers)
    rows_n = len(rows) + 1
    tbl = slide.shapes.add_table(rows_n, cols, x, y, w, h).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    # header
    for i, htxt in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = htxt
        r.font.name = FONT_BODY; r.font.size = Pt(font_size + 0.5)
        r.font.bold = True; r.font.color.rgb = WHITE
    # body
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG if (zebra and ri % 2 == 1) else WHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.name = FONT_BODY; r.font.size = Pt(font_size)
            r.font.color.rgb = INK


# ==================================================================
TOTAL = 18
# ==================================================================

# --- Slide 1: Cover ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(3.2), SW, Inches(0.04), ACCENT)
add_text(s, Inches(0.8), Inches(0.8), Inches(6), Inches(0.5),
         "BUSINESS PLAN · IR", size=13, bold=True, color=GOLD)
add_text(s, Inches(0.8), Inches(1.5), Inches(12), Inches(1.3),
         "HanView", size=80, bold=True, color=WHITE, font=FONT_TITLE)
add_text(s, Inches(0.8), Inches(2.45), Inches(12), Inches(0.7),
         "대한민국 공문서의 운영체제(OS)가 된다",
         size=26, bold=True, color=MIST, font=FONT_TITLE)
add_text(s, Inches(0.8), Inches(3.5), Inches(11), Inches(2.2),
         ("HWPX 네이티브 파서 · 시맨틱 표 이해 · AI 문서 생성 엔진.\n"
          "한컴 독점 HWP 생태계를 웹·AI 시대로 개방하는 유일한 오픈 플랫폼.\n\n"
          "· 타깃: 1인 전문사무소 → SI·공공 · BizModel: PLG + API + Enterprise\n"
          "· 90일 내 유료 MRR 1천만원, 12개월 내 시리즈 A 가능"),
         size=14, color=MIST)
add_rect(s, Inches(0.8), Inches(6.1), Inches(3), Emu(25000), ACCENT)
add_text(s, Inches(0.8), Inches(6.25), Inches(10), Inches(0.4),
         "ISM Team · license@ism-team.com · 2026.04", size=12, color=GOLD)
add_text(s, Inches(0.8), Inches(6.6), Inches(10), Inches(0.4),
         "Strictly Confidential · For Qualified Investors Only",
         size=10, color=SLATE)

# --- Slide 2: Executive Summary ---
s = prs.slides.add_slide(blank)
header(s, 2, TOTAL, "Executive Summary",
       "왜 지금, 왜 HanView인가 — 한 장으로 보는 투자 논거",
       "01 · Overview")
y = Inches(2.35)
kpi_card(s, Inches(0.5), y, Inches(3), Inches(1.2), "11조원+", "국내 HWP 기반 문서 처리 시장 (TAM)")
kpi_card(s, Inches(3.7), y, Inches(3), Inches(1.2), "0개", "HWPX 네이티브 + AI 표 이해 경쟁사", accent=GREEN)
kpi_card(s, Inches(6.9), y, Inches(3), Inches(1.2), "₩2억 → 50억", "12개월 ARR 목표 (Base case)", accent=GOLD)
kpi_card(s, Inches(10.1), y, Inches(2.7), Inches(1.2), "70%", "사내 sLM으로 AI 원가 절감")
bullet_box(s, Inches(0.5), Inches(3.85), Inches(6.2), Inches(3.3),
           "What we do",
           ["HWPX 파서 v3.0 — 단일 패스 고성능, 세계 최초 풀 스펙 구현",
            "시맨틱 표 이해 v3.0 — SpreadsheetLLM 기반, 다단계 병합표 인식",
            "AI 문서 생성 엔진 — 멀티패스, A4 1쪽 분량 섹션 확장",
            "DOCX / XLSX / MD / PDF 내보내기 — 전 포맷 왕복 호환",
            "React 라이브러리(npm) + 웹앱 + REST API 동시 제공"])
bullet_box(s, Inches(7.0), Inches(3.85), Inches(5.9), Inches(3.3),
           "Why now",
           ["공공 AI 전환 의무화 (2026 디지털플랫폼정부 기본계획)",
            "한컴 HWP 서비스 종료 루머·가격 인상 → 고객 이탈 임박",
            "LLM 토큰가 12개월 -80% → AI 문서생성 경제성 임계점",
            "HWPX = TTA·ISO 공개 표준 → 법적 리스크 ZERO",
            "국내 LLM(EXAONE·솔라) 상용화 → 데이터 주권 이슈 해소"],
           accent=GREEN)

# --- Slide 3: Problem ---
s = prs.slides.add_slide(blank)
header(s, 3, TOTAL, "Problem — 고객이 매일 겪는 고통",
       "HWP는 한국 지식노동자 3,200만명의 일일 마찰비용",
       "02 · Problem")
rows = [
    ["공공/교육", "HWP 강제 제출, 한컴 라이선스 ₩12만/년, AI 불가", "연간 ₩8,000억 기회비용"],
    ["로펌·세무·노무", "판례·계약·세무 서식이 모두 HWP, Mac에서 열기 불가", "1인 주 4.2시간 소비"],
    ["SI·SW 개발사", "HWP 파싱 라이브러리 부재 → 건당 수천만원 외주", "연간 ₩1,200억 중복개발"],
    ["스타트업·기업", "공공 RFP 대응 시 HWP 제안서 요구, 팀 협업 불가능", "수주 포기율 23%"],
    ["학교·학원", "공문 표 양식 복잡, 교사 1인 주 6시간 행정", "전국 기준 연 1.2억 시간"],
]
table(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(3.3),
      ["세그먼트", "핵심 Pain Point", "정량 영향"], rows,
      col_widths=[2, 6, 4])
add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), MIST)
add_text(s, Inches(0.8), Inches(6.1), Inches(12), Inches(0.4),
         "💡 Insight", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.8), Inches(6.45), Inches(12), Inches(0.5),
         "현존 대체재(한컴오피스·폴라리스·Google Docs)는 '뷰어'일 뿐 — AI 기반 생성·이해·자동화 계층은 전무하다.",
         size=12, color=INK)

# --- Slide 4: Solution ---
s = prs.slides.add_slide(blank)
header(s, 4, TOTAL, "Solution — HanView 4-Layer Stack",
       "파서부터 AI까지 수직 통합된 유일한 한국어 문서 플랫폼",
       "03 · Solution")
layers = [
    ("Layer 4 · Application", "웹앱 · React 위젯 · 모바일 · Slack/Teams Bot", ACCENT),
    ("Layer 3 · Intelligence", "AI 문서생성 · 시맨틱 표 이해 · 멀티패스 확장 · OCR", GOLD),
    ("Layer 2 · Runtime", "DOCX / XLSX / MD / PDF Export · 왕복 변환 · Diff", GREEN),
    ("Layer 1 · Core", "HWPX Parser v3.0 · 단일 패스 · 풀 스펙 · Worker 병렬화", NAVY),
]
ty = Inches(2.4)
for i, (name, desc, col) in enumerate(layers):
    y = ty + Inches(i * 0.9)
    add_rect(s, Inches(0.5), y, Inches(7.5), Inches(0.8), WHITE, line=LINE)
    add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.8), col)
    add_text(s, Inches(0.8), y + Inches(0.1), Inches(3), Inches(0.3),
             name, size=12, bold=True, color=col)
    add_text(s, Inches(0.8), y + Inches(0.42), Inches(7), Inches(0.3),
             desc, size=10.5, color=INK)
bullet_box(s, Inches(8.3), Inches(2.4), Inches(4.5), Inches(4.4),
           "기술적 해자",
           ["HWPX 풀 스펙 파서 (헤더·푸터·각주·하이퍼링크·필드·번호매김 전부)",
            "시맨틱 그리드 5단계 알고리즘 (그리드→코너→분류→체인→타입추론)",
            "LLM 토큰 최적화 — 40셀 단위 자동 분할 멀티패스",
            "Worker 병렬 파싱으로 100MB 문서 3초 이내",
            "React 19 + TS + Vite — 라이브러리·웹앱 이중 빌드",
            "Supabase 기반 멀티테넌시 · Tosspayments 결제"],
           accent=GREEN)

# --- Slide 5: Market Sizing ---
s = prs.slides.add_slide(blank)
header(s, 5, TOTAL, "Market Sizing — TAM / SAM / SOM",
       "보수적으로 잡아도 3년 내 ARR 500억 가능한 그릇",
       "04 · Market")
# TAM SAM SOM visual
cx, cy = Inches(3.5), Inches(4.5)
circles = [
    (Inches(3.3), NAVY, "TAM\n₩11조", "국내 전체 문서 SW·서비스 시장"),
    (Inches(2.2), ACCENT, "SAM\n₩1.8조", "HWP 의존 지식노동 영역"),
    (Inches(1.2), GOLD, "SOM\n₩180억", "3년 내 실현가능 점유"),
]
for r, col, label, desc in circles:
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - r/2, cy - r/2, r, r)
    shp.fill.solid(); shp.fill.fore_color.rgb = col
    shp.line.fill.background()
    shp.shadow.inherit = False
for i, (r, col, label, desc) in enumerate(circles):
    ty_ = Inches(2.4 + i * 1.5)
    add_text(s, Inches(7.5), ty_, Inches(1.5), Inches(0.8),
             label.split("\n")[0], size=16, bold=True, color=col)
    add_text(s, Inches(7.5), ty_ + Inches(0.35), Inches(2), Inches(0.5),
             label.split("\n")[1], size=22, bold=True, color=NAVY, font=FONT_TITLE)
    add_text(s, Inches(9.7), ty_ + Inches(0.1), Inches(3.2), Inches(1),
             desc, size=11, color=SLATE)
add_text(s, Inches(7.5), Inches(2.15), Inches(5), Inches(0.3),
         "3년 단계적 시장 포지셔닝", size=11, bold=True, color=ACCENT)
add_rect(s, Inches(7.5), Inches(6.4), Inches(5.3), Inches(0.7), BG, line=LINE)
add_text(s, Inches(7.7), Inches(6.5), Inches(5), Inches(0.5),
         "가정: SOM 10% 침투 = ARR 180억 · EV/Sales 8x = 기업가치 ₩1,440억",
         size=10.5, color=INK)

# --- Slide 6: Target Customer Segmentation ---
s = prs.slides.add_slide(blank)
header(s, 6, TOTAL, "고객 세그먼트 · 진입 우선순위",
       "Pain × WTP × 접근성으로 본 4-Quadrant",
       "05 · Target")
segs = [
    ("🎯 Beachhead", "세무·노무·행정사 1인 사무소", "12만개 사업자 · WTP ₩30K/월", ACCENT),
    ("🚀 Early Scale", "중소 로펌·회계법인", "9,000개 · WTP ₩300K/월", GOLD),
    ("💼 Enterprise", "SI·SW 개발사 (API)", "800개 · 연 ₩5천만~3억", GREEN),
    ("🏛 Public", "공공기관·지자체·학교", "17,000개 · 조달·혁신제품 트랙", NAVY),
]
for i, (tag, name, meta, col) in enumerate(segs):
    col_idx = i % 2; row_idx = i // 2
    x = Inches(0.5 + col_idx * 6.35)
    y = Inches(2.35 + row_idx * 2.4)
    add_rect(s, x, y, Inches(6.1), Inches(2.2), WHITE, line=LINE)
    add_rect(s, x, y, Inches(6.1), Inches(0.5), col)
    add_text(s, x + Inches(0.25), y + Inches(0.1), Inches(5.8), Inches(0.35),
             tag, size=12, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), y + Inches(0.65), Inches(5.8), Inches(0.5),
             name, size=17, bold=True, color=NAVY, font=FONT_TITLE)
    add_text(s, x + Inches(0.25), y + Inches(1.15), Inches(5.8), Inches(0.4),
             meta, size=11, color=SLATE)
    priority = ["진입 타이밍: 즉시 (Q2 2026)", "진입 타이밍: Q4 2026",
                "진입 타이밍: Q3 2026", "진입 타이밍: 2027 H1"]
    add_text(s, x + Inches(0.25), y + Inches(1.6), Inches(5.8), Inches(0.4),
             priority[i], size=10.5, bold=True, color=col)

# --- Slide 7: Beachhead Deep Dive ---
s = prs.slides.add_slide(blank)
header(s, 7, TOTAL, "Beachhead — 세무·노무사 1인 사무소",
       "좁게 시작해 확실하게 장악한다",
       "06 · Go-to-Market")
bullet_box(s, Inches(0.5), Inches(2.4), Inches(4.1), Inches(4.8),
           "왜 이 세그먼트인가",
           ["12만개 활성 사업자 (한국 전체 전문직 사무소 중 가장 큼)",
            "의사결정자 = 사용자 본인 → 영업사이클 2주",
            "HWP 의존도 극단적 (국세청·고용부 양식 100% HWP)",
            "월 ₩30K WTP 검증됨 (세무사랑·위멤버스 SaaS 선례)",
            "커뮤니티 결속 강함 → 입소문 바이럴 계수 1.4 예상",
            "LTV 추정 ₩720K (2년 유지) · CAC 목표 ₩80K → LTV/CAC 9x"],
           accent=ACCENT)
bullet_box(s, Inches(4.8), Inches(2.4), Inches(4.1), Inches(4.8),
           "제공 가치 (Jobs-to-be-Done)",
           ["세무조정계산서 AI 초안 생성 (4시간 → 20분)",
            "고용보험·4대보험 공문 자동 완성",
            "고객사 제출용 보고서 HWP ↔ PDF 왕복",
            "판례·국세청 예규 HWP 인덱싱·검색",
            "Mac/iPad에서 HWP 100% 열람·편집",
            "카카오톡 봇 — 사진 찍으면 공문 초안"],
           accent=GOLD)
bullet_box(s, Inches(9.1), Inches(2.4), Inches(3.7), Inches(4.8),
           "확보 채널",
           ["한국세무사회·노무사회 제휴 (회원할인)",
            "세무사랑·더존 플러그인 형태",
            "세무사 유튜버 스폰서 (택스톡·세무훈남)",
            "네이버 카페 '세무사카페' 파트너십",
            "무료 7일 + 첫달 ₩9,900 트라이얼",
            "Referral: 추천 1건당 1개월 무료"],
           accent=GREEN)

# --- Slide 8: Business Model & Pricing ---
s = prs.slides.add_slide(blank)
header(s, 8, TOTAL, "Business Model · 4-Tier Pricing",
       "PLG 깔때기로 Free → Pro → Team → Enterprise",
       "07 · Business Model")
rows = [
    ["Free", "개인", "₩0", "HWPX 뷰어·기본 변환·월 5회 AI", "SEO · 신뢰 · 퍼널 입구"],
    ["Pro", "1인 전문직", "₩9,900/월", "무제한 변환·AI 월 100회·Mac 앱", "핵심 수익원 (목표 60%)"],
    ["Team", "5~50인 법인", "₩39,000/인/월", "협업·버전관리·API 1만콜·SSO", "ACV ₩500만 이상"],
    ["Enterprise", "대기업·공공", "연 ₩5천만~3억", "온프레미스·전담 sLM·SLA·GS인증", "고마진 · 레퍼런스"],
    ["API", "개발사", "₩50/건 종량", "파서·변환·AI 엔드포인트 5종", "생태계 확장"],
]
table(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(3.0),
      ["Plan", "대상", "가격", "핵심 기능", "전략적 역할"], rows,
      col_widths=[1.3, 1.8, 1.8, 5.2, 2.4])
add_rect(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.5), BG, line=LINE)
add_text(s, Inches(0.8), Inches(5.85), Inches(12), Inches(0.4),
         "Unit Economics (Year 2 목표)", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.8), Inches(6.25), Inches(12), Inches(0.9),
         ("• Gross Margin 82% (sLM 이관 후) · CAC Payback 4.2개월 · Net Revenue Retention 118%\n"
          "• Pro ARPU ₩9.9K · Team ARPU ₩180K · Enterprise ARPU ₩12M · Blended ARPU ₩52K/월\n"
          "• 2027년 유료 계정 12,000 · MRR ₩6.2억 · ARR ₩74억 (Base)"),
         size=11, color=INK)

# --- Slide 9: Growth Flywheel ---
s = prs.slides.add_slide(blank)
header(s, 9, TOTAL, "Growth Flywheel — PLG + OSS + Community",
       "CAC=0 엔진 3종을 맞물려 돌린다",
       "08 · Growth")
steps = [
    ("① 무료 웹뷰어", "hwpx.io 도메인으로\n'hwp 열기' 키워드 SEO 독점", NAVY),
    ("② 오픈소스 파서", "Core Parser MIT 공개\nGitHub ★ 5K, 개발자 유입", ACCENT),
    ("③ 커뮤니티 템플릿", "사용자 템플릿 공유\n양식 빅데이터 축적", GOLD),
    ("④ AI 품질 격차", "학습된 템플릿으로\n경쟁사 대비 2x 품질", GREEN),
    ("⑤ 유료 전환", "무료→Pro 전환율 4.5%\nViral K=1.4", ACCENT),
]
for i, (t, d, col) in enumerate(steps):
    x = Inches(0.5 + i * 2.56)
    y = Inches(2.6)
    add_rect(s, x, y, Inches(2.4), Inches(2.5), WHITE, line=LINE)
    add_rect(s, x, y, Inches(2.4), Inches(0.45), col)
    add_text(s, x + Inches(0.15), y + Inches(0.07), Inches(2.2), Inches(0.3),
             t, size=12, bold=True, color=WHITE)
    add_text(s, x + Inches(0.15), y + Inches(0.6), Inches(2.2), Inches(1.8),
             d, size=10.5, color=INK)
    if i < 4:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   x + Inches(2.4), y + Inches(1.0),
                                   Inches(0.16), Inches(0.3))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = SLATE
        arrow.line.fill.background()
bullet_box(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.8),
           "Flywheel KPI — 90일 목표",
           ["hwpx.io MAU 50,000 (구글 SERP 1위 달성) · Pro 전환율 4.5%",
            "GitHub ★ 500 · npm weekly downloads 8,000 · Discord 멤버 1,200",
            "사용자 생성 템플릿 3,000개 · AI 품질 벤치마크 GPT-4 대비 +22%"])

# --- Slide 10: Competition & Moat ---
s = prs.slides.add_slide(blank)
header(s, 10, TOTAL, "Competitive Landscape",
       "경쟁 지형 — 2D Map으로 본 Blue Ocean",
       "09 · Competition")
rows = [
    ["한컴오피스", "HWP 편집 표준", "Mac 불가·AI 부재·비쌈", "대체 가능"],
    ["폴라리스 오피스", "저가형 뷰어", "표 깨짐·AI 부재·파싱 오류", "열세"],
    ["Google Docs", "협업 강점", "HWP 지원 0%·한국어 양식 부재", "HWP 영역 불가"],
    ["MS Copilot", "범용 AI", "HWPX 네이티브 파싱 불가", "우회 불가"],
    ["한컴독스", "HWP 웹 에디터", "AI 없음·느림·상용 API 부재", "기술 격차"],
    ["🎯 HanView", "HWPX 파서 + AI + 오픈소스", "—", "Category of One"],
]
table(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(3.5),
      ["플레이어", "강점", "약점", "당사 대비"], rows,
      col_widths=[2, 4, 4.5, 2.3])
add_rect(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.1), MIST)
add_text(s, Inches(0.8), Inches(6.25), Inches(12), Inches(0.4),
         "🛡 해자 (Moat)", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.8), Inches(6.6), Inches(12), Inches(0.6),
         ("1) 풀스펙 HWPX 파서 (2년 선행) · 2) 시맨틱 표 학습 데이터 (네트워크 효과) · "
          "3) 커뮤니티/생태계 · 4) 한국어 특화 sLM 파인튜닝"),
         size=11, color=INK)

# --- Slide 11: Product Roadmap ---
s = prs.slides.add_slide(blank)
header(s, 11, TOTAL, "Product Roadmap — 12개월 추가 개발 계획",
       "PMF → Scale → Platform 3단계",
       "10 · Roadmap")
rows = [
    ["Q2 2026", "PMF 확보", "Mac 네이티브 앱 · 카카오톡 봇 · Slack 연동 · 템플릿 마켓 MVP",
     "Beachhead 고객 100명"],
    ["Q3 2026", "AI 고도화",
     "자체 sLM 파인튜닝(EXAONE 기반) · 판례·국세청 RAG · OCR HWP 역변환 · 음성 입력",
     "토큰비 70% 절감"],
    ["Q4 2026", "협업 플랫폼",
     "실시간 공동편집 · 버전관리(Git-like) · 전자서명 · 감사로그 · Zapier 연동",
     "Team 플랜 런칭"],
    ["Q1 2027", "API 생태계",
     "Developer Portal · Webhooks · Python/Go/Java SDK · Terraform Provider",
     "개발자 5,000명"],
    ["Q2 2027", "Enterprise",
     "온프레미스 설치형 · Active Directory · DLP/DRM · GS인증 · 공공 SaaS 등록",
     "ARR ₩30억"],
    ["Q3 2027+", "글로벌",
     "다국어 확장 (日 ichitaro, 中 WPS), PDF Native Edit, Agent 자동화",
     "일본 파일럿"],
]
table(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(4.5),
      ["분기", "테마", "핵심 개발 항목", "달성 KPI"], rows,
      col_widths=[1.3, 1.8, 6.7, 2.5], font_size=9.5)

# --- Slide 12: Tech Architecture ---
s = prs.slides.add_slide(blank)
header(s, 12, TOTAL, "Tech Architecture · 추가 인프라 계획",
       "현재 스택 + 12개월 보강 계획",
       "11 · Tech")
bullet_box(s, Inches(0.5), Inches(2.4), Inches(6.1), Inches(4.7),
           "현재 (v4.0 as-is)",
           ["Frontend: React 19 · Vite 7 · TS 5.9 · Zustand",
            "Parser: HWPX v3.0 · Web Worker 병렬 · 단일 패스",
            "AI: OpenAI API · 시맨틱 표 처리 · 멀티패스 확장",
            "Export: docx · exceljs · jspdf · html2canvas",
            "Auth/DB: Supabase · RLS 멀티테넌시",
            "Payment: Tosspayments 구독",
            "배포: Docker · Vercel · SSR 부분"],
           accent=SLATE)
bullet_box(s, Inches(6.7), Inches(2.4), Inches(6.1), Inches(4.7),
           "Next 12개월 보강",
           ["sLM: EXAONE 3.5 · vLLM · GPU A100×2 (KT Cloud)",
            "Vector DB: pgvector 확장 · 템플릿·판례 임베딩 500만건",
            "Observability: OpenTelemetry · Grafana · Sentry",
            "보안: SOC2 Type I · ISMS-P · 개인정보영향평가",
            "CDN: Cloudflare R2 (문서 저장) · 한국 리전",
            "API: Kong Gateway · Rate limit · Circuit Breaker",
            "Mobile: React Native (iOS/Android) · Tauri (Mac)"],
           accent=ACCENT)

# --- Slide 13: 90-Day Execution Plan ---
s = prs.slides.add_slide(blank)
header(s, 13, TOTAL, "90-Day Execution Plan",
       "실행 가능한 주간 단위 액션 플랜",
       "12 · Execution")
rows = [
    ["Week 1-2", "법인 설립·브랜드", "hwpx.io 도메인 · 법인 전환 · 상표출원", "법인등기"],
    ["Week 3-4", "무료 뷰어 런칭", "랜딩 SEO · 30개 키워드 상위 10위", "MAU 5K"],
    ["Week 5-6", "오픈소스 공개", "Core Parser MIT 릴리즈 · HN/PH 론칭", "GitHub ★ 500"],
    ["Week 7-8", "Pro 플랜 오픈", "Stripe·Toss 결제 · 트라이얼 설계", "유료 30"],
    ["Week 9-10", "세무사 50 베타", "세무사회 공동 세미나 · 케이스 10건", "MRR ₩1M"],
    ["Week 11-12", "SI 파트너", "NHN·LG CNS·삼성SDS 3곳 미팅 · LOI 2건", "LOI 2"],
    ["Week 13", "시드 IR", "투자자 15팀 미팅 · Term Sheet", "₩15억 Term"],
]
table(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(4.3),
      ["기간", "마일스톤", "주요 액션", "성공 지표"], rows,
      col_widths=[1.5, 2.5, 6, 2.3], font_size=10.5)
add_rect(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.35), ACCENT)
add_text(s, Inches(0.8), Inches(6.95), Inches(12), Inches(0.3),
         "책임자: CEO(영업·투자) · CTO(개발) · CMO(마케팅·SEO) — 주간 OKR 점검",
         size=10.5, bold=True, color=WHITE)

# --- Slide 14: Financials ---
s = prs.slides.add_slide(blank)
header(s, 14, TOTAL, "Financial Projection — 3-Year",
       "보수적 Base Case · 월별 코호트 기반 롤업",
       "13 · Financials")
rows = [
    ["유료 계정 (누적)", "200", "3,500", "12,000", "28,000"],
    ["MRR (₩M)", "3", "85", "620", "1,850"],
    ["ARR (₩억)", "0.4", "10", "74", "222"],
    ["Gross Margin", "55%", "72%", "82%", "85%"],
    ["S&M 비중", "85%", "55%", "38%", "28%"],
    ["월 버닝 (₩M)", "35", "95", "120", "180"],
    ["Runway 소진 시점", "—", "-", "흑자전환", "순이익 ₩45억"],
    ["헤드카운트", "6", "18", "42", "78"],
]
table(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(3.8),
      ["지표", "2026 Q2", "2026 EOY", "2027 EOY", "2028 EOY"], rows,
      col_widths=[4, 2, 2, 2, 2], font_size=11)
add_rect(s, Inches(0.5), Inches(6.4), Inches(6), Inches(0.8), BG, line=LINE)
add_text(s, Inches(0.7), Inches(6.5), Inches(5.7), Inches(0.3),
         "Base Case 가정", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(6.8), Inches(5.7), Inches(0.4),
         "월 Churn 3% · Net Add 900/월 · ARPU ₩52K",
         size=10, color=INK)
add_rect(s, Inches(6.8), Inches(6.4), Inches(6), Inches(0.8), MIST, line=LINE)
add_text(s, Inches(7), Inches(6.5), Inches(5.7), Inches(0.3),
         "Exit 가상 시나리오", size=11, bold=True, color=GREEN)
add_text(s, Inches(7), Inches(6.8), Inches(5.7), Inches(0.4),
         "2028 ARR ₩220억 × EV/Sales 8x = 밸류 ₩1,760억",
         size=10, color=INK)

# --- Slide 15: Funding Ask ---
s = prs.slides.add_slide(blank)
header(s, 15, TOTAL, "Funding Ask — Seed ₩15억",
       "18개월 러닝 + 시리즈 A 준비까지 커버",
       "14 · Funding")
# Pie-like use of funds
uses = [
    ("제품 개발 (45%)", "₩6.75억", "sLM 파인튜닝·모바일·엔터프라이즈 보안", ACCENT),
    ("인재 채용 (30%)", "₩4.5억", "백엔드 3·AI 2·세일즈 2·CS 1", GOLD),
    ("마케팅 (15%)", "₩2.25억", "SEO·콘텐츠·세무사회 파트너", GREEN),
    ("인증·운영 (10%)", "₩1.5억", "ISMS-P·GS인증·법무·운영비", NAVY),
]
for i, (name, amt, desc, col) in enumerate(uses):
    y = Inches(2.4 + i * 0.95)
    add_rect(s, Inches(0.5), y, Inches(7.8), Inches(0.85), WHITE, line=LINE)
    add_rect(s, Inches(0.5), y, Inches(0.12), Inches(0.85), col)
    add_text(s, Inches(0.75), y + Inches(0.08), Inches(3), Inches(0.35),
             name, size=13, bold=True, color=NAVY)
    add_text(s, Inches(0.75), y + Inches(0.45), Inches(6), Inches(0.35),
             desc, size=10.5, color=SLATE)
    add_text(s, Inches(6.3), y + Inches(0.2), Inches(1.8), Inches(0.5),
             amt, size=16, bold=True, color=col,
             align=PP_ALIGN.RIGHT, font=FONT_TITLE)
bullet_box(s, Inches(8.5), Inches(2.4), Inches(4.3), Inches(4.3),
           "투자 조건",
           ["라운드: Seed", "규모: ₩15억", "밸류에이션: Pre ₩60억 (Post ₩75억)",
            "형태: CPS 또는 SAFE (post-money)",
            "Use of Fund 기간: 18개월",
            "시리즈 A 목표: 2027 Q3, ₩80억@₩400억"],
           accent=ACCENT)
add_rect(s, Inches(8.5), Inches(6.9), Inches(4.3), Inches(0.3), ACCENT)
add_text(s, Inches(8.7), Inches(6.93), Inches(4), Inches(0.3),
         "Lead 1곳 + Co-investor 2곳 희망",
         size=10, bold=True, color=WHITE)

# --- Slide 16: Risks ---
s = prs.slides.add_slide(blank)
header(s, 16, TOTAL, "Risks & Mitigation",
       "예상 리스크와 선제 대응 전략",
       "15 · Risks")
rows = [
    ["한컴의 법적 대응", "중", "상", "HWPX는 TTA·ISO 표준. 변호사 자문 완료. 공격적 커뮤니케이션"],
    ["AI 토큰비 폭증", "상", "중", "사내 sLM 이관 Q3 2026 · 캐시·배치 최적화"],
    ["핵심 인력 이탈", "중", "상", "ESOP 15% 풀 · 기술리드 3인 2년 lock-up"],
    ["공공조달 지연", "상", "중", "PLG 현금흐름으로 의존도 최소화. 혁신제품 병행"],
    ["대기업 복제 진입", "하", "상", "오픈소스·커뮤니티 해자. 파트너십으로 선점"],
    ["개인정보·보안 이슈", "중", "상", "ISMS-P 2026 취득. 온프레미스 옵션. SOC2 로드맵"],
    ["환율·거시 경기", "중", "중", "국내 매출 90%로 환리스크 낮음. 고정비 slim"],
]
table(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(4.5),
      ["리스크", "발생 확률", "영향도", "대응 전략"], rows,
      col_widths=[3, 1.5, 1.5, 6.3], font_size=10.5)

# --- Slide 17: Team ---
s = prs.slides.add_slide(blank)
header(s, 17, TOTAL, "Team & Hiring Plan",
       "핵심 창업팀 + 18개월 핵심 영입 포지션",
       "16 · Team")
bullet_box(s, Inches(0.5), Inches(2.4), Inches(6.1), Inches(4.7),
           "현재 팀 (ISM Team)",
           ["CEO — 풀스택·HWPX 도메인 10년 · 제품 총괄",
            "CTO — 컴파일러·파서 전문 · Parser v3 아키텍트",
            "AI Lead — LLM·시맨틱 그리드 · 논문 3편",
            "FE Engineer — React 19·라이브러리 패키징",
            "Designer — 제품 UX · 브랜딩",
            "BD — 공공·법인 영업 10년 경력",
            "—— 총 6명 · 평균 경력 8년"],
           accent=NAVY)
bullet_box(s, Inches(6.7), Inches(2.4), Inches(6.1), Inches(4.7),
           "Next 12 Hires (Seed 자금 기반)",
           ["1. VP of Sales (엔터프라이즈 판매 헤드)",
            "2. Growth PM (PLG·SEO·퍼널 최적화)",
            "3-4. Backend Engineer × 2 (Go·분산시스템)",
            "5-6. AI/ML Engineer × 2 (sLM·RAG·vLLM)",
            "7. Mobile Engineer (React Native·Tauri)",
            "8. Security Engineer (ISMS-P·SOC2)",
            "9-10. Account Executive × 2 (세무·법무 세그먼트)",
            "11. Customer Success Manager",
            "12. Content/Community Manager (개발자 생태계)"],
           accent=ACCENT)

# --- Slide 18: Closing ---
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(3.2), SW, Inches(0.04), ACCENT)
add_text(s, Inches(0.8), Inches(0.8), Inches(6), Inches(0.5),
         "THE ASK", size=13, bold=True, color=GOLD)
add_text(s, Inches(0.8), Inches(1.5), Inches(12), Inches(1.3),
         "함께 가시겠습니까?", size=60, bold=True, color=WHITE, font=FONT_TITLE)
add_text(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(2.5),
         ("HanView는 한국 지식노동의 마찰을 제거하는 플랫폼입니다.\n\n"
          "· Seed ₩15억, Pre-money ₩60억\n"
          "· 18개월 내 ARR ₩74억, Series A Ready\n"
          "· '공공·법률·세무'를 장악하고 일본·아시아로 확장\n\n"
          "HWP는 가두고 있고, 우리는 열 것입니다."),
         size=16, color=MIST)
add_rect(s, Inches(0.8), Inches(6.1), Inches(3), Emu(25000), ACCENT)
add_text(s, Inches(0.8), Inches(6.25), Inches(12), Inches(0.4),
         "Contact · license@ism-team.com · hanview.ai",
         size=14, bold=True, color=GOLD)
add_text(s, Inches(0.8), Inches(6.65), Inches(12), Inches(0.4),
         "© 2026 ISM Team · Strictly Confidential",
         size=10, color=SLATE)

# ---- Save ----
out = "/Users/k1/Documents/Project/HanView_React/hanview-react-app/docs/HanView_Business_Plan_2026Q2.pptx"
prs.save(out)
print(f"SAVED: {out}")
